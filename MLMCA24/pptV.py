from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
presentation = Presentation()

# Slide 1: Title Slide
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[0]

title.text = "Commonly Used Visualization Charts and Plots in Python"
subtitle.text = "Understanding which chart to use for different problems"

# Slide 2: Introduction
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[0]

title.text = "Introduction"
content.text = (
    "Briefly introduce the importance of data visualization in analysis.\n"
    "Mention that the right choice of charts enhances communication and understanding of data."
)

# Define slide content for each chart
chart_data = [
    ("Line Chart", "Ideal for showing trends or changes over time.", "Plotting stock prices over months."),
    ("Bar Chart", "Suitable for comparing categories or showing the distribution of a single variable.", "Comparing sales performance of different products."),
    ("Histogram", "Visualizes the distribution of a single variable.", "Displaying the frequency distribution of exam scores."),
    ("Scatter Plot", "Depicts the relationship between two numerical variables.", "Examining the correlation between hours of study and exam scores."),
    ("Pie Chart", "Represents parts of a whole and is useful for displaying proportions.", "Showing the percentage distribution of expenses in a budget."),
    ("Box Plot", "Displays the distribution of a dataset and highlights outliers.", "Analyzing the distribution of salaries in a company."),
    ("Heatmap", "Represents data values in a matrix format using color intensity.", "Displaying correlations in a correlation matrix."),
    ("Area Chart", "Shows the magnitude of a phenomenon over time.", "Displaying the growth of a company's revenue over years."),
]

# Create slides for each chart
for chart_title, chart_description, chart_example in chart_data:
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[0]
    
    title.text = chart_title
    content.text = f"{chart_description}\nExample: {chart_example}"

# Slide 11: Conclusion
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[0]

title.text = "Conclusion"
content.text = (
    "Summarize key points.\n"
    "Emphasize the importance of choosing the right visualization for effective communication."
)

# Slide 12: Q&A
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[0]

title.text = "Q&A"
content.text = "Invite questions and discussions."

# Save the presentation
presentation.save("VisualizationChartsPresentation.pptx")
